"""
readers.py — one uniform view over four file formats.

Every reader turns a file into an ordered list of BLOCKS:

    {"type": "paragraph", "text": str,  "location": str, "runs": [...]}
    {"type": "table",     "rows": [[str, ...], ...], "location": str}

Because all four formats produce the same shape, the rule engine is written
once and works everywhere. The `location` field is what makes citations
accurate per format, and it differs deliberately by format:

    .docx  -> the nearest preceding heading, e.g. "§B.4 Project Team"
              (Word has no stored page concept - pagination is computed at
              render time - so a section anchor is the honest citation and is
              also more useful: a reviewer can Ctrl+F it.)
    .pdf   -> "p. 4"   (PDF really does have pages)
    .pptx  -> "Slide 7"
    .xlsx  -> "Budget Summary!P10" / the sheet name

Format reliability note: .docx and .xlsx expose real table structure, so
cell-based rules on them are exact. PDF has NO native concept of a table -
pdfplumber infers tables from ruled lines and alignment - so the same rule is
inherently less reliable there. The UI surfaces this rather than hiding it.
"""

import os

from docx import Document
from docx.document import Document as _DocumentClass
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
import openpyxl


# --------------------------------------------------------------------------
# Shared helpers
# --------------------------------------------------------------------------

def _para_block(text, location, runs=None):
    return {"type": "paragraph", "text": text, "location": location,
            "runs": runs or []}


def _table_block(rows, location):
    return {"type": "table", "rows": rows, "location": location}


# --------------------------------------------------------------------------
# WORD
# --------------------------------------------------------------------------

def iter_block_items(parent):
    """
    Yield each paragraph and table IN TRUE DOCUMENT ORDER.

    python-docx exposes .paragraphs and .tables as two separate lists and
    silently loses how they were interleaved. Walking the raw XML body restores
    the real order, which is what makes verbatim section capture faithful
    (paragraph -> table -> paragraph stays in that order).
    """
    parent_elm = parent.element.body if isinstance(parent, _DocumentClass) else parent._tc
    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def _distinct_cells(row):
    """Row cells with horizontal-merge duplicates removed.

    Word repeats the SAME underlying cell across a horizontal merge. Without
    de-duplicating, "the cell to the right of the label" can land on a repeat
    of the label's own cell and return the wrong value.
    """
    out, seen = [], set()
    for cell in row.cells:
        key = id(cell._tc)
        if key not in seen:
            seen.add(key)
            out.append(cell)
    return out


class DocxReader:
    kind = "docx"

    def __init__(self, path):
        self.path = path
        self.doc = Document(path)

    def blocks(self):
        out = []
        current = "body"
        for b in iter_block_items(self.doc):
            if isinstance(b, Paragraph):
                style = ""
                try:
                    style = (b.style.name or "") if b.style else ""
                except Exception:
                    style = ""
                if b.text.strip() and "heading" in style.lower():
                    current = "§" + b.text.strip()
                runs = [{"text": r.text, "bold": bool(r.bold), "italic": bool(r.italic)}
                        for r in b.runs if r.text]
                out.append(_para_block(b.text, current, runs))
            else:
                rows = [[c.text.strip() for c in _distinct_cells(r)] for r in b.rows]
                out.append(_table_block(rows, current))
        return out


# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------

class PdfReader:
    kind = "pdf"

    def __init__(self, path):
        self.path = path

    def blocks(self):
        import pdfplumber
        out = []
        with pdfplumber.open(self.path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                loc = f"p. {i}"
                try:
                    for t in page.extract_tables():
                        rows = [[(c or "").strip() for c in row] for row in t]
                        if rows:
                            out.append(_table_block(rows, loc))
                except Exception:
                    pass  # borderless/odd tables: skip rather than fail the run
                text = page.extract_text() or ""
                for line in text.split("\n"):
                    if line.strip():
                        out.append(_para_block(line.strip(), loc))
        return out


# --------------------------------------------------------------------------
# POWERPOINT
# --------------------------------------------------------------------------

class PptxReader:
    kind = "pptx"

    def __init__(self, path):
        self.path = path

    def blocks(self):
        from pptx import Presentation
        out = []
        prs = Presentation(self.path)
        for i, slide in enumerate(prs.slides, start=1):
            loc = f"Slide {i}"
            for shape in slide.shapes:
                if shape.has_table:
                    rows = [[c.text.strip() for c in r.cells] for r in shape.table.rows]
                    out.append(_table_block(rows, loc))
                elif shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = "".join(r.text for r in para.runs) or para.text
                        if txt.strip():
                            out.append(_para_block(txt.strip(), loc))
            # Speaker notes are often where the real detail lives.
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        out.append(_para_block(notes.strip(), f"{loc} (notes)"))
            except Exception:
                pass
        return out


# --------------------------------------------------------------------------
# EXCEL
# --------------------------------------------------------------------------

class XlsxReader:
    kind = "xlsx"

    def __init__(self, path):
        self.path = path
        # data_only=True returns the value Excel last SAVED, not formula text.
        self.wb = openpyxl.load_workbook(path, data_only=True)

    def blocks(self):
        """Each sheet becomes one table block (bounded, so huge sheets stay sane)."""
        out = []
        for name in self.wb.sheetnames:
            ws = self.wb[name]
            rows = []
            for row in ws.iter_rows(max_row=min(ws.max_row or 1, 400),
                                    max_col=min(ws.max_column or 1, 60)):
                rows.append([("" if c.value is None else str(c.value)) for c in row])
            if rows:
                out.append(_table_block(rows, name))
        return out

    def cell_range(self, sheet_name, ref):
        """
        Read a rectangular range like "P10:R14" and return it as a grid of
        strings (rows of cells), so it can be rendered as a table.
        Returns (rows, problem).
        """
        target = None
        want = (sheet_name or "").strip().lower()
        for n in self.wb.sheetnames:
            if n.strip().lower() == want:
                target = self.wb[n]
                break
        if target is None:
            return None, (f"No tab named '{sheet_name}'. "
                          f"Tabs found: {', '.join(self.wb.sheetnames)}.")
        try:
            cells = target[ref]
        except Exception as exc:
            return None, f"Could not read range {ref}: {exc}"
        if not isinstance(cells, tuple):
            cells = ((cells,),)
        rows = []
        for row in cells:
            if not isinstance(row, tuple):
                row = (row,)
            rows.append([("" if c.value is None else str(c.value)) for c in row])
        if not any(any(c.strip() for c in r) for r in rows):
            return None, (f"Range {ref} on '{sheet_name}' is empty, or holds "
                          f"formulas whose results were never saved by Excel.")
        return rows, None

    def cell(self, sheet_name, cell_ref):
        """
        Return (value, problem). openpyxl cannot calculate formulas; it can only
        return what Excel cached on last save. A workbook generated by a program
        and never opened in Excel may therefore have no value at all - that is a
        distinct failure state, not an empty cell, and is reported as such.
        """
        target = None
        want = (sheet_name or "").strip().lower()
        for n in self.wb.sheetnames:
            if n.strip().lower() == want:
                target = self.wb[n]
                break
        if target is None:
            return None, (f"No tab named '{sheet_name}'. "
                          f"Tabs found: {', '.join(self.wb.sheetnames)}.")
        try:
            value = target[cell_ref].value
        except Exception as exc:
            return None, f"Could not read cell {cell_ref}: {exc}"
        if value is None:
            return None, (f"Cell {cell_ref} on '{sheet_name}' is empty, or holds a "
                          f"formula whose result was never saved by Excel. Open the "
                          f"workbook in Excel, let it recalculate, save, and retry.")
        return str(value), None


# --------------------------------------------------------------------------
# Dispatch
# --------------------------------------------------------------------------

READERS = {".docx": DocxReader, ".pdf": PdfReader,
           ".pptx": PptxReader, ".xlsx": XlsxReader, ".xlsm": XlsxReader}

SUPPORTED = sorted(READERS.keys())


def open_reader(path):
    """Return the right reader for a file, or (None, problem-message)."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".doc":
        return None, ("Legacy '.doc' files cannot be read. Open in Word and use "
                      "File > Save As > Word Document (.docx), then re-upload.")
    if ext not in READERS:
        return None, f"Unsupported file type '{ext}'. Supported: {', '.join(SUPPORTED)}."
    try:
        return READERS[ext](path), None
    except Exception as exc:
        return None, f"Could not open this file: {exc}"
